"""
Artifact Generator - Produce output files (DOCX, XLSX, PDF, PPTX).

Phase 6 integration for generating document artifacts from structured data.

Architecture Reference:
- architecture/concepts/ARTIFACT_SYSTEM.md
- architecture/BENCHMARK_ALIGNMENT.md (output artifact generation)
"""

import json
import logging
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, TYPE_CHECKING

if TYPE_CHECKING:
    from libs.gateway.persistence.turn_manager import TurnDirectory

logger = logging.getLogger(__name__)


@dataclass
class Artifact:
    """A generated output artifact."""
    artifact_id: str
    artifact_type: str  # docx, xlsx, pdf, pptx
    filename: str
    path: str
    size_bytes: int
    title: str = ""
    description: str = ""
    created_at: str = ""

    def to_dict(self) -> Dict[str, Any]:
        return {
            "artifact_id": self.artifact_id,
            "type": self.artifact_type,
            "filename": self.filename,
            "path": self.path,
            "size_bytes": self.size_bytes,
            "title": self.title,
            "description": self.description,
            "created_at": self.created_at,
        }


@dataclass
class ArtifactManifest:
    """Manifest tracking all artifacts for a turn."""
    turn_id: str
    artifacts: List[Artifact] = field(default_factory=list)
    created_at: str = ""

    def to_dict(self) -> Dict[str, Any]:
        return {
            "turn_id": self.turn_id,
            "artifacts": [a.to_dict() for a in self.artifacts],
            "total_artifacts": len(self.artifacts),
            "created_at": self.created_at,
        }

    def add(self, artifact: Artifact) -> None:
        self.artifacts.append(artifact)

    def get_references(self) -> str:
        """Get markdown references for synthesis."""
        if not self.artifacts:
            return ""
        lines = ["### Generated Artifacts", ""]
        for a in self.artifacts:
            lines.append(f"- **{a.title or a.filename}** ({a.artifact_type.upper()}, {a.size_bytes} bytes): `{a.filename}`")
        return "\n".join(lines)


class DocxGenerator:
    """Generate DOCX documents."""

    def generate(
        self,
        output_path: Path,
        title: str,
        content: str,
        sections: Optional[List[Dict[str, str]]] = None,
        metadata: Optional[Dict[str, str]] = None,
    ) -> Artifact:
        """
        Generate a DOCX file.

        Args:
            output_path: Path to write the file
            title: Document title
            content: Main content (markdown-like)
            sections: Optional list of {heading, body} dicts
            metadata: Optional metadata (author, subject, etc.)

        Returns:
            Artifact describing the generated file
        """
        lines = []

        # Build document content as plain text (fallback if python-docx unavailable)
        lines.append(f"# {title}")
        lines.append("")
        if metadata:
            for key, value in metadata.items():
                lines.append(f"**{key}:** {value}")
            lines.append("")

        if content:
            lines.append(content)
            lines.append("")

        if sections:
            for section in sections:
                heading = section.get("heading", "")
                body = section.get("body", "")
                lines.append(f"## {heading}")
                lines.append("")
                lines.append(body)
                lines.append("")

        text_content = "\n".join(lines)

        try:
            from docx import Document
            doc = Document()
            doc.add_heading(title, 0)

            if metadata:
                for key, value in metadata.items():
                    doc.add_paragraph(f"{key}: {value}")

            if content:
                for para in content.split("\n\n"):
                    doc.add_paragraph(para.strip())

            if sections:
                for section in sections:
                    doc.add_heading(section.get("heading", ""), level=1)
                    for para in section.get("body", "").split("\n\n"):
                        doc.add_paragraph(para.strip())

            doc.save(str(output_path))
        except ImportError:
            # Fallback: write as .docx-compatible text
            output_path.write_text(text_content)

        size = output_path.stat().st_size

        return Artifact(
            artifact_id=f"docx_{datetime.now().strftime('%Y%m%d%H%M%S')}",
            artifact_type="docx",
            filename=output_path.name,
            path=str(output_path),
            size_bytes=size,
            title=title,
            description=f"Document: {title}",
            created_at=datetime.now().isoformat(),
        )


class XlsxGenerator:
    """Generate XLSX spreadsheets."""

    def generate(
        self,
        output_path: Path,
        title: str,
        sheets: List[Dict[str, Any]],
    ) -> Artifact:
        """
        Generate an XLSX file.

        Args:
            output_path: Path to write the file
            title: Spreadsheet title
            sheets: List of sheet definitions with name, headers, rows

        Returns:
            Artifact describing the generated file
        """
        try:
            import openpyxl
            wb = openpyxl.Workbook()

            for i, sheet_def in enumerate(sheets):
                if i == 0:
                    ws = wb.active
                    ws.title = sheet_def.get("name", "Sheet1")
                else:
                    ws = wb.create_sheet(sheet_def.get("name", f"Sheet{i+1}"))

                headers = sheet_def.get("headers", [])
                rows = sheet_def.get("rows", [])

                if headers:
                    ws.append(headers)
                for row in rows:
                    ws.append(row)

            wb.save(str(output_path))
        except ImportError:
            # Fallback: write as CSV
            lines = []
            for sheet_def in sheets:
                lines.append(f"--- {sheet_def.get('name', 'Sheet')} ---")
                headers = sheet_def.get("headers", [])
                rows = sheet_def.get("rows", [])
                if headers:
                    lines.append(",".join(str(h) for h in headers))
                for row in rows:
                    lines.append(",".join(str(c) for c in row))
                lines.append("")
            output_path.write_text("\n".join(lines))

        size = output_path.stat().st_size

        return Artifact(
            artifact_id=f"xlsx_{datetime.now().strftime('%Y%m%d%H%M%S')}",
            artifact_type="xlsx",
            filename=output_path.name,
            path=str(output_path),
            size_bytes=size,
            title=title,
            description=f"Spreadsheet: {title}",
            created_at=datetime.now().isoformat(),
        )


class PdfGenerator:
    """Generate PDF documents."""

    def generate(
        self,
        output_path: Path,
        title: str,
        content: str,
        sections: Optional[List[Dict[str, str]]] = None,
    ) -> Artifact:
        """
        Generate a PDF file.

        Args:
            output_path: Path to write the file
            title: Document title
            content: Main content
            sections: Optional sections

        Returns:
            Artifact describing the generated file
        """
        # Build text content
        lines = [title, "=" * len(title), "", content, ""]
        if sections:
            for section in sections:
                lines.append(section.get("heading", ""))
                lines.append("-" * len(section.get("heading", "")))
                lines.append(section.get("body", ""))
                lines.append("")

        text_content = "\n".join(lines)

        try:
            import fitz  # PyMuPDF
            doc = fitz.open()
            page = doc.new_page()

            # Simple text rendering
            text_point = fitz.Point(72, 72)
            page.insert_text(text_point, text_content[:3000], fontsize=11)

            doc.save(str(output_path))
            doc.close()
        except ImportError:
            # Fallback: write as text
            output_path.write_text(text_content)

        size = output_path.stat().st_size

        return Artifact(
            artifact_id=f"pdf_{datetime.now().strftime('%Y%m%d%H%M%S')}",
            artifact_type="pdf",
            filename=output_path.name,
            path=str(output_path),
            size_bytes=size,
            title=title,
            description=f"PDF: {title}",
            created_at=datetime.now().isoformat(),
        )


class PptxGenerator:
    """Generate PPTX presentations."""

    def generate(
        self,
        output_path: Path,
        title: str,
        slides: List[Dict[str, Any]],
    ) -> Artifact:
        """
        Generate a PPTX file.

        Args:
            output_path: Path to write the file
            title: Presentation title
            slides: List of slide definitions with title, content, notes

        Returns:
            Artifact describing the generated file
        """
        try:
            from pptx import Presentation
            prs = Presentation()

            for slide_def in slides:
                layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_def.get("title", "")
                if slide.placeholders[1]:
                    slide.placeholders[1].text = slide_def.get("content", "")

            prs.save(str(output_path))
        except ImportError:
            # Fallback: write as text
            lines = [f"Presentation: {title}", ""]
            for i, slide_def in enumerate(slides):
                lines.append(f"--- Slide {i+1}: {slide_def.get('title', '')} ---")
                lines.append(slide_def.get("content", ""))
                if slide_def.get("notes"):
                    lines.append(f"Notes: {slide_def['notes']}")
                lines.append("")
            output_path.write_text("\n".join(lines))

        size = output_path.stat().st_size

        return Artifact(
            artifact_id=f"pptx_{datetime.now().strftime('%Y%m%d%H%M%S')}",
            artifact_type="pptx",
            filename=output_path.name,
            path=str(output_path),
            size_bytes=size,
            title=title,
            description=f"Presentation: {title}",
            created_at=datetime.now().isoformat(),
        )


class ArtifactGenerator:
    """
    Unified artifact generator for Phase 6.

    Dispatches to type-specific generators and maintains manifest.
    """

    def __init__(self, output_dir: Optional[Path] = None):
        self.output_dir = output_dir or Path("artifacts")
        self.docx = DocxGenerator()
        self.xlsx = XlsxGenerator()
        self.pdf = PdfGenerator()
        self.pptx = PptxGenerator()
        self.manifest: Optional[ArtifactManifest] = None

    def create_manifest(self, turn_id: str) -> ArtifactManifest:
        """Create a new manifest for a turn."""
        self.manifest = ArtifactManifest(
            turn_id=turn_id,
            created_at=datetime.now().isoformat(),
        )
        return self.manifest

    def generate(
        self,
        artifact_type: str,
        filename: str,
        **kwargs,
    ) -> Artifact:
        """
        Generate an artifact of the specified type.

        Args:
            artifact_type: Type (docx, xlsx, pdf, pptx)
            filename: Output filename
            **kwargs: Type-specific arguments

        Returns:
            Generated Artifact
        """
        self.output_dir.mkdir(parents=True, exist_ok=True)
        output_path = self.output_dir / filename

        if artifact_type == "docx":
            artifact = self.docx.generate(output_path, **kwargs)
        elif artifact_type == "xlsx":
            artifact = self.xlsx.generate(output_path, **kwargs)
        elif artifact_type == "pdf":
            artifact = self.pdf.generate(output_path, **kwargs)
        elif artifact_type == "pptx":
            artifact = self.pptx.generate(output_path, **kwargs)
        else:
            raise ValueError(f"Unknown artifact type: {artifact_type}")

        # Add to manifest
        if self.manifest:
            self.manifest.add(artifact)

        return artifact

    def save_manifest(self, turn_dir: "TurnDirectory") -> Path:
        """Save artifact manifest to turn directory."""
        if not self.manifest:
            self.manifest = ArtifactManifest(
                turn_id="unknown",
                created_at=datetime.now().isoformat(),
            )

        manifest_path = turn_dir.doc_path("artifact_manifest.json")
        manifest_path.write_text(json.dumps(self.manifest.to_dict(), indent=2))
        return manifest_path

    @staticmethod
    def load_manifest(turn_dir: "TurnDirectory") -> Optional[ArtifactManifest]:
        """Load artifact manifest from turn directory."""
        manifest_path = turn_dir.doc_path("artifact_manifest.json")
        if not manifest_path.exists():
            return None

        try:
            data = json.loads(manifest_path.read_text())
            manifest = ArtifactManifest(
                turn_id=data.get("turn_id", ""),
                created_at=data.get("created_at", ""),
            )
            for a in data.get("artifacts", []):
                manifest.add(Artifact(
                    artifact_id=a.get("artifact_id", ""),
                    artifact_type=a.get("type", ""),
                    filename=a.get("filename", ""),
                    path=a.get("path", ""),
                    size_bytes=a.get("size_bytes", 0),
                    title=a.get("title", ""),
                    description=a.get("description", ""),
                    created_at=a.get("created_at", ""),
                ))
            return manifest
        except Exception:
            return None
